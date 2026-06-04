"""
Generate the Creative Matrix Activity Guide PowerPoint.
Run: python3 generate_activity_guide.py
Output: Creative Matrix - Activity Guide.pptx (in the same directory)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Colours ──────────────────────────────────────────────────────────────────
BLACK      = RGBColor(0x16, 0x16, 0x16)
WHITE      = RGBColor(0xFC, 0xFC, 0xFC)
ORANGE     = RGBColor(0xFF, 0x4E, 0x00)
DARK_GRAY  = RGBColor(0x48, 0x48, 0x48)
MED_GRAY   = RGBColor(0x7A, 0x7A, 0x7A)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

# ── Dimensions (match reference: 13.33" × 7.50") ─────────────────────────────
W = Emu(12192000)
H = Emu(6858000)

# ── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # blank
    return prs.slides.add_slide(layout)

def bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))

def add_para(tf, text, size, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             space_before=0, space_after=0, level=0):
    """Append a paragraph to a text frame."""
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_before = Pt(space_before)
    p.space_after  = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size  = _Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return p

def first_para(tf, text, size, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    """Populate the first (default) paragraph of a text frame."""
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return p

def label(slide, text, x, y, w=5.41, color=ORANGE):
    """Small all-caps section label (e.g. 'STEP 1', 'A.I. PROMPT')."""
    tb = box(slide, x, y, w, 0.30)
    tb.text_frame.word_wrap = False
    first_para(tb.text_frame, text, 9, bold=True, color=color)
    return tb

def section_tag(slide, text):
    """Top-left section identifier used on content slides."""
    tb = box(slide, 0.66, 0.43, 6.0, 0.61)
    tb.text_frame.word_wrap = False
    first_para(tb.text_frame, text, 14, bold=False, color=MED_GRAY)
    return tb

def add_rect(slide, l, t, w, h, fill_color):
    """Add a filled rectangle (no text)."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

# ── Slide builders ────────────────────────────────────────────────────────────

def make_cover(prs):
    slide = blank_slide(prs)
    bg(slide, BLACK)

    # Orange accent line
    add_rect(slide, 0.60, 0.55, 0.06, 3.80, ORANGE)

    # Big headline
    tb = box(slide, 0.80, 0.60, 10.5, 3.80)
    tf = tb.text_frame
    tf.word_wrap = True
    first_para(tf, "Construct a Creative Matrix", 52, bold=True, color=WHITE)

    # Method name
    tb2 = box(slide, 0.80, 4.85, 6.0, 0.40)
    first_para(tb2.text_frame, "Creative Matrix", 14, color=MED_GRAY)

    # Label
    tb3 = box(slide, 0.80, 5.30, 4.0, 0.35)
    first_para(tb3.text_frame, "Activity guide", 12, color=ORANGE)

    return slide


def make_toc(prs):
    slide = blank_slide(prs)
    bg(slide, WHITE)
    section_tag(slide, "Table of contents")

    items = [
        ("01", "What is this activity"),
        ("02", "How to do it"),
        ("03", "Tips for success"),
        ("04", "Examples"),
    ]
    col_x = [0.66, 3.68, 6.83, 10.03]

    for i, ((num, title), x) in enumerate(zip(items, col_x)):
        tb = box(slide, x, 5.40, 2.85, 0.65)
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = num + "  "; r1.font.size = Pt(18); r1.font.bold = True; r1.font.color.rgb = ORANGE
        r2 = p.add_run(); r2.text = title; r2.font.size = Pt(18); r2.font.bold = False; r2.font.color.rgb = BLACK

    return slide


def make_what_is_this(prs):
    slide = blank_slide(prs)
    bg(slide, WHITE)
    section_tag(slide, "What is this activity?")

    # Left column — headline + body
    tb = box(slide, 0.67, 1.50, 5.80, 5.60)
    tf = tb.text_frame
    tf.word_wrap = True

    first_para(tf, "DESIGN THE MATRIX THAT GENERATES IDEAS WORTH ACTING ON", 20,
               bold=True, color=BLACK)

    add_para(tf, "", 6)  # spacer

    body = (
        "Creative Matrix generates ideas by forcing unlikely combinations — "
        "pairing specific challenges against categories of activators to produce "
        "ideas that wouldn't emerge from open brainstorming. The quality of your "
        "matrix determines the quality of your ideas: the right columns and the "
        "right enablers produce output that is grounded in real insights and "
        "surprising enough to be useful.\n\n"
        "This guide walks you through constructing a matrix that is ready to run "
        "in person, on a shared whiteboard, or as a self-advancing PowerPoint deck."
    )
    add_para(tf, body, 14, color=DARK_GRAY)

    add_para(tf, "", 8)

    add_para(tf, "WHAT YOU WILL NEED", 11, bold=True, color=BLACK)
    needs = [
        "Your topic or challenge — a clear overarching problem or opportunity",
        "Upstream work — research notes, insights, personas, or a journey map",
        "Access to an AI tool (Claude, ChatGPT, or Copilot) for optional assistance",
        "At least 30 minutes to work through all steps",
    ]
    for n in needs:
        add_para(tf, "·  " + n, 13, color=DARK_GRAY)

    return slide


def make_how_to_do_it(prs):
    slide = blank_slide(prs)
    bg(slide, WHITE)
    section_tag(slide, "How to do it")

    tb = box(slide, 0.67, 1.60, 8.0, 5.20)
    tf = tb.text_frame
    tf.word_wrap = True

    steps = [
        ("Step 1:", "Define your challenge", "5 min"),
        ("Step 2:", "Build your columns", "15 min"),
        ("Step 3:", "Select your enablers", "10 min"),
        ("Step 4:", "Add think-about prompts", "10 min"),
        ("Step 5:", "Choose your delivery format", "5 min"),
    ]
    first_para(tf, "", 4)  # top spacer
    for num, name, timing in steps:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        r1 = p.add_run(); r1.text = num + "  "; r1.font.size = Pt(20); r1.font.bold = True; r1.font.color.rgb = ORANGE
        r2 = p.add_run(); r2.text = name; r2.font.size = Pt(20); r2.font.bold = False; r2.font.color.rgb = BLACK
        r3 = p.add_run(); r3.text = "  (" + timing + ")"; r3.font.size = Pt(16); r3.font.color.rgb = MED_GRAY

    return slide


def make_step(prs, step_num, step_name, timing, instructions,
              right_label, right_content, right_is_prompt=False):
    """
    Standard two-column step slide.
    instructions: list of strings (bullet points)
    right_content: list of strings for the right panel
    right_is_prompt: if True, style right panel as AI prompt box
    """
    slide = blank_slide(prs)
    bg(slide, WHITE)

    # Left: step number tag
    label(slide, f"STEP {step_num}", 0.59, 0.24, color=ORANGE)

    # Left: step name
    tb_name = box(slide, 0.67, 0.65, 5.33, 1.21)
    tf = tb_name.text_frame
    tf.word_wrap = True
    first_para(tf, step_name, 24, bold=True, color=BLACK)

    # Left: instructions
    tb_inst = box(slide, 0.67, 2.10, 5.33, 4.00)
    tf2 = tb_inst.text_frame
    tf2.word_wrap = True
    first_para(tf2, "", 4)
    for instr in instructions:
        add_para(tf2, "·  " + instr, 14, color=DARK_GRAY, space_before=4)

    # Left: timing (bottom)
    tb_time = box(slide, 0.67, 6.46, 2.0, 0.30)
    first_para(tb_time.text_frame, timing, 12, color=MED_GRAY)

    # Right: label
    right_label_color = ORANGE if right_is_prompt else MED_GRAY
    label(slide, right_label, 7.54, 0.24, w=5.41, color=right_label_color)

    # Right: content panel
    if right_is_prompt:
        # Tinted background box
        add_rect(slide, 7.40, 0.55, 5.50, 6.65, LIGHT_GRAY)

    tb_right = box(slide, 7.54, 0.60, 5.20, 6.55)
    tf3 = tb_right.text_frame
    tf3.word_wrap = True
    first_para(tf3, "", 4)
    for line in right_content:
        size = 12 if right_is_prompt else 13
        color = DARK_GRAY
        add_para(tf3, line, size, color=color, space_before=3)

    return slide


def make_tips(prs, tips):
    slide = blank_slide(prs)
    bg(slide, WHITE)
    section_tag(slide, "Tips for success")

    tb = box(slide, 0.67, 1.20, 11.8, 5.80)
    tf = tb.text_frame
    tf.word_wrap = True
    first_para(tf, "", 4)
    for tip in tips:
        add_para(tf, "·  " + tip, 14, color=DARK_GRAY, space_before=6)

    return slide


def make_pitfalls(prs, pitfalls):
    """pitfalls: list of (bad, good) tuples"""
    slide = blank_slide(prs)
    bg(slide, WHITE)
    section_tag(slide, "Common pitfalls to avoid")

    tb = box(slide, 0.67, 1.20, 11.8, 5.80)
    tf = tb.text_frame
    tf.word_wrap = True
    first_para(tf, "", 4)
    for bad, good in pitfalls:
        p_bad = tf.add_paragraph()
        p_bad.space_before = Pt(8)
        r = p_bad.add_run(); r.text = "❌  "; r.font.size = Pt(14); r.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)
        r2 = p_bad.add_run(); r2.text = bad; r2.font.size = Pt(14); r2.font.color.rgb = DARK_GRAY

        p_good = tf.add_paragraph()
        p_good.space_before = Pt(2)
        p_good.space_after = Pt(4)
        r3 = p_good.add_run(); r3.text = "✅  "; r3.font.size = Pt(14); r3.font.color.rgb = RGBColor(0x00, 0x80, 0x00)
        r4 = p_good.add_run(); r4.text = good; r4.font.size = Pt(14); r4.font.color.rgb = DARK_GRAY

    return slide


def make_examples_divider(prs, program_name="[Program]"):
    slide = blank_slide(prs)
    bg(slide, BLACK)

    tb = box(slide, 0.66, 2.60, 10.39, 1.20)
    tf = tb.text_frame
    tf.word_wrap = True
    first_para(tf, program_name + " examples", 40, bold=False, color=WHITE,
               align=PP_ALIGN.LEFT)

    tb2 = box(slide, 0.66, 3.90, 10.39, 0.50)
    first_para(tb2.text_frame, "Creative Matrix", 18, color=ORANGE)

    tb3 = box(slide, 0.66, 4.50, 10.39, 0.40)
    first_para(tb3.text_frame, "Examples from other sessions", 14, color=MED_GRAY)

    return slide


def make_example_placeholder(prs, num):
    slide = blank_slide(prs)
    bg(slide, LIGHT_GRAY)
    tb = box(slide, 0.66, 3.0, 11.8, 1.0)
    first_para(tb.text_frame,
               f"[Example slide {num} — populate with a completed Creative Matrix example]",
               18, color=MED_GRAY, align=PP_ALIGN.CENTER)
    return slide


# ── Content ───────────────────────────────────────────────────────────────────

STEP1_INSTR = [
    "Write your topic at the top — a clear phrase naming the problem or opportunity you're exploring (e.g. 'Improving first-year employee onboarding' or 'Driving adoption of our analytics platform')",
    "Below it, write your overarching challenge — the broadest How Might We statement this session will address",
    "Be specific: vague topics produce generic ideas; a sharply framed challenge produces ideas you can act on",
]
STEP1_RIGHT = [
    "OVERARCHING TOPIC",
    "",
    "[Write your topic here]",
    "",
    "──────────────────────────────────",
    "",
    "OVERARCHING CHALLENGE",
    "",
    "How might we...",
    "",
    "[Write your challenge statement here]",
]

STEP2_INSTR = [
    "Choose your column type based on the upstream work you have:",
    "   · How Might We statements — use when you have research insights or strategic priorities to turn into opportunity prompts",
    "   · Customer segments or stakeholder roles — use when distinct groups need meaningfully different ideas",
    "   · Experience stages — use when you're designing a service or journey (default: Entice, Enter, Engage, Exit, Extend)",
    "Generate four column headings of your chosen type",
    "Use the A.I. prompt on the right if you have research notes or context to feed in",
]
STEP2_RIGHT_PROMPT = [
    "I'm constructing a Creative Matrix for the following",
    "challenge:",
    "",
    "Topic: [INSERT YOUR TOPIC]",
    "",
    "Context: [DESCRIBE YOUR ORGANISATION AND",
    "AUDIENCE — e.g. 'a hospital improving patient intake'",
    "or 'a tech company improving tool adoption']",
    "",
    "Upstream work:",
    "[PASTE RESEARCH NOTES, INSIGHT THEMES, AFFINITY",
    "CLUSTERS, STAKEHOLDER MAP, OR JOURNEY MAP]",
    "",
    "Please:",
    "1. Recommend the best column type (HMW statements,",
    "   customer segments, or experience stages) and explain why",
    "2. Generate four column headings of that type, grounded",
    "   in the upstream work I've provided",
]

STEP3_INSTR = [
    "Start with the default four enablers: Technology & Digital Media, Games & Competitions, People & Partnerships, Events & Programs",
    "Adjust for context — add something currently relevant (AI, a key platform, a strategic priority) or swap in experience-specific enablers if your columns are journey stages",
    "Choose enablers that feel slightly uncomfortable — tension between column and enabler produces surprising ideas",
    "The bottom row is always Wild Card — leave it open for ideas that don't fit a named enabler",
    "Use the A.I. prompt on the right to get context-specific recommendations",
]
STEP3_RIGHT_PROMPT = [
    "I'm building a Creative Matrix on the following topic:",
    "",
    "Topic: [INSERT TOPIC]",
    "",
    "Columns: [LIST YOUR FOUR COLUMN HEADINGS]",
    "",
    "Context: [DESCRIBE YOUR ORGANISATION AND AUDIENCE]",
    "",
    "Please:",
    "1. Recommend four enablers that will create the most",
    "   productive tension with these columns",
    "2. For each enabler, provide 4–6 sub-bullet 'think about'",
    "   prompts — specific examples within that category that",
    "   would spark ideas given this particular context",
    "3. Briefly explain why each enabler creates useful tension",
    "   for this topic and challenge set",
]

STEP4_INSTR = [
    "Under each enabler heading, list 4–6 specific sub-examples within that category",
    "These are the 'think about' prompts that help people ideate when they get stuck — 'Technology & Digital Media' is abstract; 'wearables, social media, generative AI, embedded sensors' gives people something to grab onto",
    "Make the prompts specific to your topic — generic sub-bullets produce generic ideas",
    "If you used the A.I. prompt in Step 3, your sub-bullets are already populated — review and refine them",
]
STEP4_RIGHT = [
    "EXAMPLE",
    "",
    "Enabler: Games & Competitions",
    "",
    "Think about:",
    "  · Motivations and what drives behaviour",
    "  · Rewards, badges, points, prizes",
    "  · Scoring and leaderboards",
    "  · Teamwork and cooperative challenges",
    "  · Peer competition and accountability",
    "",
    "──────────────────────────────────",
    "",
    "Now write sub-bullets for each of",
    "your four enablers in the same format.",
]

STEP5_INSTR = [
    "Your completed matrix can be run in three ways — choose based on your team's setup:",
    "   · Physical wall: print and mount the grid; use index cards for column and row headers; participants use sticky notes",
    "   · Miro or Mural: pre-build the grid on a shared board; participants add digital sticky notes simultaneously",
    "   · Auto-advance PowerPoint deck: build the 16-slide session deck using your matrix; participants work individually against slides that advance every minute",
    "For this program, use the auto-advance PowerPoint format — take your completed matrix spec and use the separate Auto-Advance Deck template to build the session-ready file",
]
STEP5_RIGHT = [
    "YOUR COMPLETED MATRIX SPEC",
    "",
    "Topic:",
    "[Your topic]",
    "",
    "Overarching challenge:",
    "[Your HMW statement]",
    "",
    "Column 1:  [Heading]",
    "Column 2:  [Heading]",
    "Column 3:  [Heading]",
    "Column 4:  [Heading]",
    "",
    "Enabler 1:  [Name + sub-bullets]",
    "Enabler 2:  [Name + sub-bullets]",
    "Enabler 3:  [Name + sub-bullets]",
    "Enabler 4:  [Name + sub-bullets]",
]

TIPS = [
    "Define your challenge sharply before building the matrix — a vague topic produces generic ideas; a specific challenge produces ideas you can act on",
    "Choose your column type based on your upstream work, not habit — if you have research insights, use HMW statements; if you have distinct stakeholder groups, use segments; if you're designing a journey, use experience stages",
    "Pick enablers that make you slightly uncomfortable — the tension is productive; obvious enablers produce obvious ideas",
    "Populate the sub-bullets before the session — four to six specific examples per enabler make the difference between participants who get stuck and participants who fill the grid in eight minutes",
    "The bottom row is always Wild Card — don't skip it; good ideas that don't fit named enablers need somewhere to go",
    "If you don't have upstream research, use AI to generate column headings from your topic — then treat those as hypotheses to validate with real people before investing in them",
]

PITFALLS = [
    (
        "Column headings that are all variations of the same thing ('improve the experience', 'improve the service', 'improve the product')",
        "Make each column distinct — different dimensions of the challenge, not rephrasing of one idea. If columns blend together, the matrix produces redundant ideas."
    ),
    (
        "Enablers that are too safe or familiar (e.g. 'email', 'meetings', 'training sessions')",
        "Choose enablers that push the team somewhere they wouldn't naturally go — if the enabler feels obvious, replace it with something that creates productive discomfort."
    ),
    (
        "Skipping the sub-bullets because the enabler name seems self-explanatory",
        "Populate 4–6 specific sub-examples per enabler — they do the priming work that makes the sprint fast and generative, especially for participants who aren't natural ideators."
    ),
    (
        "Building the matrix during the session while participants wait",
        "Construct the full matrix — columns, enablers, and sub-bullets — before anyone arrives. Construction is preparation; the session is ideation."
    ),
    (
        "Choosing all four default enablers without considering context",
        "Start with the default four, then ask: is there something currently de rigueur (AI, a key platform, a strategic priority) that belongs here instead? The best enabler sets are tuned to the moment."
    ),
]


# ── Build ─────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()

    make_cover(prs)
    make_toc(prs)
    make_what_is_this(prs)
    make_how_to_do_it(prs)

    make_step(prs, 1, "Define your challenge", "5 minutes",
              STEP1_INSTR, "YOUR CHALLENGE", STEP1_RIGHT)

    make_step(prs, 2, "Build your columns", "15 minutes",
              STEP2_INSTR, "A.I. PROMPT", STEP2_RIGHT_PROMPT, right_is_prompt=True)

    make_step(prs, 3, "Select your enablers", "10 minutes",
              STEP3_INSTR, "A.I. PROMPT", STEP3_RIGHT_PROMPT, right_is_prompt=True)

    make_step(prs, 4, "Add think-about prompts", "10 minutes",
              STEP4_INSTR, "EXAMPLE", STEP4_RIGHT)

    make_step(prs, 5, "Choose your delivery format", "5 minutes",
              STEP5_INSTR, "YOUR MATRIX SPEC", STEP5_RIGHT)

    make_tips(prs, TIPS)
    make_pitfalls(prs, PITFALLS)
    make_examples_divider(prs)
    make_example_placeholder(prs, 1)
    make_example_placeholder(prs, 2)

    out = os.path.join(os.path.dirname(__file__), "Creative Matrix - Activity Guide.pptx")
    prs.save(out)
    print(f"Saved: {out}")

if __name__ == "__main__":
    build()
